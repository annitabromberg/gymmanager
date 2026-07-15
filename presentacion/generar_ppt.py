from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Tema
bg = RGBColor(12, 22, 43)
accent = RGBColor(244, 114, 182)
card = RGBColor(24, 38, 59)
text = RGBColor(248, 250, 252)
muted = RGBColor(203, 213, 225)

# ---------- Portada ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg

shape = slide.shapes.add_shape(1, 0, 0, 13.333, 7.5)
shape.fill.solid()
shape.fill.fore_color.rgb = bg
shape.line.color.rgb = bg

# Título
box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(12), Inches(1.2))
text_frame = box.text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
p.text = "GymManager"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = accent

box2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10), Inches(1.2))
text_frame2 = box2.text_frame
text_frame2.clear()
p2 = text_frame2.paragraphs[0]
p2.text = "Presentación para Fundamentos de Informática"
p2.font.size = Pt(18)
p2.font.color.rgb = text

box3 = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(10), Inches(1.2))
text_frame3 = box3.text_frame
text_frame3.clear()
p3 = text_frame3.paragraphs[0]
p3.text = "Anna Bromberg | María Josefina Velázquez Paz | Ángela Romero"
p3.font.size = Pt(16)
p3.font.color.rgb = muted

# Línea decorativa
line = slide.shapes.add_shape(1, Inches(0.8), Inches(4.7), Inches(5.2), Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = accent
line.line.color.rgb = accent

# ----- Diapositiva 1 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg

slide.shapes.title.text = "1. ¿De qué se trata la aplicación?"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "GymManager es una aplicación web para gestionar un gimnasio.",
    "Permite registrar alumnos, rutinas, asistencia y cuotas.",
    "La idea es organizar la información del gimnasio en una sola plataforma.",
    "La aplicación fue desarrollada como proyecto para Fundamentos de Informática.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = text
    p.level = 0

# ----- Diapositiva 2 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg
slide.shapes.title.text = "2. ¿Qué herramientas se usaron?"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "Python: lenguaje principal para la lógica.",
    "Flask: framework para construir páginas web.",
    "HTML: estructura de las páginas.",
    "CSS: diseño y apariencia moderna.",
    "JSON: almacenamiento del registro de alumnos y otras listas.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = text

# ----- Diapositiva 3 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg
slide.shapes.title.text = "3. ¿Cómo se organizó el proyecto?"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "app.py: inicia la aplicación.",
    "templates/: contiene las páginas HTML.",
    "static/: contiene CSS, videos e imágenes.",
    "gymmanager_app/: guarda la lógica y las rutas de la app.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = text

# ----- Diapositiva 4 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg
slide.shapes.title.text = "4. Estructura del código"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "app.py: prende la aplicación y define el puerto.",
    "gymmanager_app/__init__.py: crea la instancia principal de Flask.",
    "gymmanager_app/routes.py: define rutas como login, inicio, alumnos y asistencia.",
    "gymmanager_app/data.py: lee y guarda información en archivos JSON.",
    "Todo está separado por tareas para que sea más fácil de entender y mantener.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = text

# ----- Diapositiva 5 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg
slide.shapes.title.text = "5. Cómo se conecta todo"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "1) El usuario abre una URL en el navegador.",
    "2) Flask reconoce la ruta pedida.",
    "3) La ruta llama a una función del código.",
    "4) Esa función consulta o guarda información y devuelve una página HTML.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = text

# ----- Diapositiva 6 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg
slide.shapes.title.text = "6. Capturas de pantalla de la interfaz"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "Se pueden incluir capturas del login, la pantalla de inicio y la lista de alumnos.",
    "Estas imágenes ayudan a que la presentación sea más visual y fácil de entender.",
    "Además, muestran cómo se ve la app desde la perspectiva del usuario final.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = text

# ----- Diapositiva 7 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg
slide.shapes.title.text = "7. Paso a paso: cómo funciona"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "1) El usuario entra a la aplicación.",
    "2) Se muestra el login.",
    "3) Si las credenciales son correctas, accede al sistema.",
    "4) Puede gestionar alumnos, asistencia y rutinas.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = text

# ----- Diapositiva 8 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg
slide.shapes.title.text = "5. ¿Qué pasa con los datos?"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "La información se guarda en archivos JSON.",
    "alumnos.json: alumnos registrados.",
    "rutinas.json: planes de entrenamiento.",
    "asistencia.json: seguimiento de presencia.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = text

# ----- Diapositiva 6 -----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg
slide.shapes.title.text = "6. Conclusión"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, line in enumerate([
    "GymManager es un ejemplo sencillo de cómo una idea real puede transformarse en una aplicación web.",
    "Se combina lógica, diseño y organización de información.",
    "Es una buena forma de comprender cómo se construyen proyectos de software.",
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = text

prs.save("presentacion/gymmanager_presentacion.pptx")
print("PowerPoint generado: presentacion/gymmanager_presentacion.pptx")
